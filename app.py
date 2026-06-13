from flask import Flask, render_template, request, jsonify, redirect, url_for, flash, session
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from werkzeug.security import generate_password_hash, check_password_hash
from datetime import datetime
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import plotly.utils
import json, io, os

# ── Optional AI libs ──────────────────────────────────────────────────────────
try:
    from google import genai as genai_new
    GEMINI_LIB = True
    GEMINI_NEW_SDK = True
except ImportError:
    GEMINI_NEW_SDK = False
    try:
        import google.generativeai as genai_old
        GEMINI_LIB = True
    except ImportError:
        GEMINI_LIB = False

# Gemini via direct HTTP — works without any SDK
import urllib.request
import urllib.error

def call_gemini_http(api_key, prompt):
    """Call Gemini API using only built-in Python urllib — no extra packages needed."""
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-pro-latest:generateContent?key={api_key}"
    body = json.dumps({
        "contents": [{"parts": [{"text": prompt}]}]
    }).encode('utf-8')
    req = urllib.request.Request(url, data=body,
                                  headers={'Content-Type': 'application/json'}, method='POST')
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            result = json.loads(resp.read().decode('utf-8'))
            return result['candidates'][0]['content']['parts'][0]['text']
    except urllib.error.HTTPError as e:
        err = e.read().decode('utf-8')
        raise RuntimeError(f"Gemini API error: {err}")
    except Exception as e:
        raise RuntimeError(f"Gemini request failed: {e}")
try:
    from openai import OpenAI; OPENAI_LIB = True
except: OPENAI_LIB = False
try:
    import anthropic; CLAUDE_LIB = True
except: CLAUDE_LIB = False

# ── Optional file libs ────────────────────────────────────────────────────────
try:
    import pdfplumber; PDF_OK = True
except: PDF_OK = False
try:
    from docx import Document as DocxDoc; DOCX_OK = True
except: DOCX_OK = False
try:
    from pptx import Presentation; PPTX_OK = True
except: PPTX_OK = False

# ─────────────────────────────────────────────────────────────────────────────
#  PASTE YOUR API KEYS HERE
#  Gemini  (FREE) → https://aistudio.google.com/app/apikey
#  DeepSeek(cheap)→ https://platform.deepseek.com/api_keys
#  ChatGPT (paid) → https://platform.openai.com/api-keys
#  Claude  (paid) → https://console.anthropic.com/settings/keys
# ─────────────────────────────────────────────────────────────────────────────
API_KEYS = {
    "gemini":   os.environ.get('GEMINI_API_KEY',   'YOUR_GEMINI_API_KEY'),
    "chatgpt":  os.environ.get('OPENAI_API_KEY',   'YOUR_OPENAI_API_KEY'),
    "claude":   os.environ.get('CLAUDE_API_KEY',   'YOUR_CLAUDE_API_KEY'),
    "deepseek": os.environ.get('DEEPSEEK_API_KEY', 'YOUR_DEEPSEEK_API_KEY'),
}

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TMP_DIR  = os.path.join(BASE_DIR, 'tmp_data')
os.makedirs(TMP_DIR, exist_ok=True)

app = Flask(__name__)
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'biz-analyst-secret-2024')
app.config['SQLALCHEMY_DATABASE_URI'] = os.environ.get(
    'DATABASE_URL',
    f"sqlite:///{os.path.join(BASE_DIR, 'bizanalyst.db')}"
)
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024

db = SQLAlchemy(app)
login_manager = LoginManager(app)
login_manager.login_view = 'login'

COLORS = ['#89b4fa','#cba6f7','#a6e3a1','#f38ba8','#fab387','#f9e2af','#74c7ec','#94e2d5']
ALLOWED = {'csv','xlsx','xls','ods','pdf','docx','doc','txt','rtf','pptx','ppt',
           'json','xml','yaml','yml','db','sqlite','sqlite3','png','jpg','jpeg','html','htm','log','md'}

THEMES = {
    'dark':     {'bg':'#0d0d1a','surface':'#13131f','card':'#1a1a2e','border':'#252540','accent':'#89b4fa','accent2':'#cba6f7','text':'#cdd6f4'},
    'midnight': {'bg':'#0a0a0f','surface':'#0f0f1a','card':'#141428','border':'#1e1e38','accent':'#c084fc','accent2':'#e879f9','text':'#e2e8f0'},
    'forest':   {'bg':'#0a1a0f','surface':'#0f1f14','card':'#14281a','border':'#1e3824','accent':'#4ade80','accent2':'#86efac','text':'#dcfce7'},
    'ocean':    {'bg':'#0a1520','surface':'#0f1e2d','card':'#14273a','border':'#1e3550','accent':'#38bdf8','accent2':'#7dd3fc','text':'#e0f2fe'},
    'sunset':   {'bg':'#1a0d0a','surface':'#1f1308','card':'#281a0f','border':'#3a2410','accent':'#fb923c','accent2':'#fbbf24','text':'#fff7ed'},
}

# ── Database Models ───────────────────────────────────────────────────────────
class User(UserMixin, db.Model):
    id             = db.Column(db.Integer, primary_key=True)
    username       = db.Column(db.String(80), unique=True, nullable=False)
    email          = db.Column(db.String(120), unique=True, nullable=False)
    password_hash  = db.Column(db.String(256), nullable=False)
    role           = db.Column(db.String(20), default='user')
    theme          = db.Column(db.String(20), default='dark')
    analyses_count = db.Column(db.Integer, default=0)
    created_at     = db.Column(db.DateTime, default=datetime.utcnow)

    def set_password(self, pw):   self.password_hash = generate_password_hash(pw)
    def check_password(self, pw): return check_password_hash(self.password_hash, pw)

@login_manager.user_loader
def load_user(uid): return User.query.get(int(uid))

# ── Context processor — inject theme into every template ─────────────────────
@app.context_processor
def inject_theme():
    if current_user.is_authenticated:
        t = THEMES.get(current_user.theme, THEMES['dark'])
    else:
        t = THEMES['dark']
    return dict(theme=t)

# ── Helpers ───────────────────────────────────────────────────────────────────
def get_ext(fn): return fn.rsplit('.',1)[1].lower() if '.' in fn else ''
def key_set(p):  return API_KEYS.get(p,'') not in ('','YOUR_GEMINI_API_KEY','YOUR_OPENAI_API_KEY','YOUR_CLAUDE_API_KEY','YOUR_DEEPSEEK_API_KEY')
def tmp_path():
    uid = current_user.id if current_user.is_authenticated else 'anon'
    return os.path.join(TMP_DIR, f'data_{uid}.json')

def get_theme():
    if current_user.is_authenticated:
        return THEMES.get(current_user.theme, THEMES['dark'])
    return THEMES['dark']

# ── AI dispatcher ─────────────────────────────────────────────────────────────
def call_ai(provider, prompt):
    p = provider.lower()
    if p == 'gemini':
        if not API_KEYS.get('gemini') or API_KEYS['gemini'] == 'YOUR_GEMINI_API_KEY':
            raise RuntimeError("No Gemini API key set in app.py")
        # Always use direct HTTP — no SDK needed
        return call_gemini_http(API_KEYS['gemini'], prompt)
    if p == 'chatgpt':
        if not OPENAI_LIB: raise RuntimeError("pip install openai")
        c = OpenAI(api_key=API_KEYS['chatgpt'])
        return c.chat.completions.create(model="gpt-4o-mini",
            messages=[{"role":"user","content":prompt}],max_tokens=2000).choices[0].message.content
    if p == 'claude':
        if not CLAUDE_LIB: raise RuntimeError("pip install anthropic")
        c = anthropic.Anthropic(api_key=API_KEYS['claude'])
        return c.messages.create(model="claude-3-haiku-20240307",max_tokens=2000,
            messages=[{"role":"user","content":prompt}]).content[0].text
    if p == 'deepseek':
        if not OPENAI_LIB: raise RuntimeError("pip install openai")
        c = OpenAI(api_key=API_KEYS['deepseek'],base_url="https://api.deepseek.com")
        return c.chat.completions.create(model="deepseek-chat",
            messages=[{"role":"user","content":prompt}],max_tokens=2000).choices[0].message.content
    raise RuntimeError(f"Unknown provider: {provider}")

def get_active_provider():
    for p in ['gemini','chatgpt','claude','deepseek']:
        if key_set(p): return p
    return None

# ── File readers ──────────────────────────────────────────────────────────────
def read_csv(file):
    content = file.read()
    for enc in ['utf-8','latin-1','cp1252']:
        try: return pd.read_csv(io.BytesIO(content), encoding=enc)
        except: pass
    raise ValueError("Cannot decode CSV")

def read_excel(file):
    content = file.read()
    xl = pd.ExcelFile(io.BytesIO(content))
    if len(xl.sheet_names) == 1:
        return pd.read_excel(io.BytesIO(content), sheet_name=xl.sheet_names[0])
    frames = []
    for s in xl.sheet_names:
        try:
            df = pd.read_excel(io.BytesIO(content), sheet_name=s)
            df['_sheet'] = s; frames.append(df)
        except: pass
    return pd.concat(frames, ignore_index=True) if frames else pd.read_excel(io.BytesIO(content))

def read_pdf(file):
    if not PDF_OK: raise ValueError("pip install pdfplumber")
    content = file.read(); tables, texts = [], []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t: texts.append(t)
            for tbl in page.extract_tables():
                if tbl and len(tbl)>1: tables.append(tbl)
    for t in tables:
        try:
            df = pd.DataFrame(t[1:], columns=t[0]).dropna(how='all')
            if len(df)>0:
                for col in df.columns:
                    try: df[col]=pd.to_numeric(df[col].astype(str).str.replace(',','').str.strip())
                    except: pass
                return df
        except: continue
    return pd.DataFrame({'Content':['\n'.join(texts)],'Pages':[len(texts)]})

def read_docx(file):
    if not DOCX_OK: raise ValueError("pip install python-docx")
    content = file.read(); doc = DocxDoc(io.BytesIO(content))
    if doc.tables:
        t = doc.tables[0]
        headers = [c.text.strip() for c in t.rows[0].cells]
        rows = [[c.text.strip() for c in r.cells] for r in t.rows[1:]]
        df = pd.DataFrame(rows, columns=headers)
        for col in df.columns:
            try: df[col]=pd.to_numeric(df[col].str.replace(',','').str.strip())
            except: pass
        return df
    return pd.DataFrame({'Content':[p.text for p in doc.paragraphs if p.text.strip()]})

def read_pptx(file):
    if not PPTX_OK: raise ValueError("pip install python-pptx")
    content = file.read(); prs = Presentation(io.BytesIO(content)); data=[]
    for i,slide in enumerate(prs.slides,1):
        texts=[s.text.strip() for s in slide.shapes if hasattr(s,'text') and s.text.strip()]
        data.append({'Slide':i,'Content':' | '.join(texts)})
    return pd.DataFrame(data)

def read_json(file):
    data=json.loads(file.read())
    if isinstance(data,list): return pd.json_normalize(data)
    for v in data.values():
        if isinstance(v,list) and v: return pd.json_normalize(v)
    return pd.json_normalize([data])

def read_text(file):
    content=file.read().decode('utf-8',errors='replace')
    lines=content.strip().split('\n')
    if lines and ',' in lines[0]:
        try: return pd.read_csv(io.StringIO(content))
        except: pass
    return pd.DataFrame({'Line':range(1,len(lines)+1),'Content':lines})

def read_file(file):
    ext=get_ext(file.filename)
    readers={'csv':read_csv,'xlsx':read_excel,'xls':read_excel,'ods':read_excel,
             'pdf':read_pdf,'docx':read_docx,'doc':read_docx,'pptx':read_pptx,'ppt':read_pptx,
             'json':read_json,'txt':read_text,'rtf':read_text,'log':read_text,'md':read_text,
             'html':lambda f: pd.read_html(io.BytesIO(f.read()))[0],
             'htm':lambda f: pd.read_html(io.BytesIO(f.read()))[0],
             'xml':lambda f: pd.read_xml(io.BytesIO(f.read())),
             'yaml':lambda f: _read_yaml(f),
    }
    if ext not in readers: raise ValueError(f"Unsupported: .{ext}")
    return readers[ext](file)

def _read_yaml(file):
    try:
        import yaml
        data = yaml.safe_load(file.read())
        return pd.json_normalize(data) if isinstance(data, list) else pd.DataFrame([data])
    except ImportError:
        raise ValueError("pip install pyyaml")

# ── Chart helpers ─────────────────────────────────────────────────────────────
def base_layout(title='', theme=None):
    t = theme or THEMES['dark']
    return dict(
        paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)',
        font=dict(color=t['text'], family='Inter,sans-serif', size=12),
        margin=dict(l=50, r=30, t=55, b=50),
        title=dict(text=title, font=dict(size=14, color=t['text']), x=0.01),
        legend=dict(bgcolor='rgba(0,0,0,0)', font=dict(color=t['text']),
                    bordercolor=t['border'], borderwidth=1),
        xaxis=dict(gridcolor=t['border'], zerolinecolor=t['border'],
                   color=t['text'], showgrid=True, tickangle=-30),
        yaxis=dict(gridcolor=t['border'], zerolinecolor=t['border'],
                   color=t['text'], showgrid=True),
        hoverlabel=dict(bgcolor=t['card'], font_color=t['text'], bordercolor=t['border']),
    )

def make_chart(fig, title, icon):
    enc = plotly.utils.PlotlyJSONEncoder()
    return {'title': title, 'icon': icon, 'data': json.loads(enc.encode(fig))}

def smart_detect_date_col(df):
    """Detect date/time columns for accurate trend analysis."""
    for col in df.columns:
        if df[col].dtype == 'object':
            try:
                converted = pd.to_datetime(df[col], infer_datetime_format=True, errors='coerce')
                if converted.notna().sum() > len(df) * 0.7:
                    return col, converted
            except Exception:
                pass
    return None, None

def clean_numeric(df, col):
    """Return clean numeric series without nulls."""
    return df[col].dropna().replace([float('inf'), float('-inf')], float('nan')).dropna()

def generate_charts(df, analysis_type='general', specific_cols=None, theme=None):
    charts = []
    t = theme or THEMES['dark']
    cs = [t['accent'], t['accent2']] + COLORS

    # Smart column detection
    num = df.select_dtypes(include='number').columns.tolist()
    cat = df.select_dtypes(include='object').columns.tolist()

    # Remove low-quality columns
    num = [c for c in num if clean_numeric(df, c).nunique() > 1]
    cat = [c for c in cat if df[c].nunique() <= 50 and df[c].nunique() > 1]

    fn   = specific_cols or []
    fc   = [c for c in fn if c in num] or num[:6]
    fcat = [c for c in fn if c in cat] or cat[:3]
    a    = analysis_type.lower()

    # Detect date column
    date_col, date_series = smart_detect_date_col(df)

    def safe(fn):
        try: fn()
        except Exception: pass

    # ── 1. Sorted Bar Chart (most accurate for comparisons) ──
    if fcat and fc:
        def _bar():
            g = df.groupby(fcat[0])[fc[0]].agg(['sum','mean','count']).reset_index()
            g.columns = [fcat[0], 'Total', 'Average', 'Count']
            g = g.sort_values('Total', ascending=False).head(15)
            fig = px.bar(g, x=fcat[0], y='Total',
                         text='Total',
                         color='Total',
                         color_continuous_scale='Blues',
                         labels={'Total': fc[0], fcat[0]: fcat[0]},
                         hover_data={'Average': ':.2f', 'Count': True})
            fig.update_traces(
                texttemplate='%{text:,.0f}',
                textposition='outside',
                marker_line_width=0
            )
            fig.update_layout(**base_layout(f'Total {fc[0]} by {fcat[0]}', t))
            fig.update_layout(showlegend=False)
            charts.append(make_chart(fig, f'Total {fc[0]} by {fcat[0]}', '📊'))
        safe(_bar)

    # ── 2. Accurate Donut with percentages ──
    if fcat:
        def _pie():
            vc = df[fcat[0]].value_counts().head(8).reset_index()
            vc.columns = [fcat[0], 'Count']
            vc['Percentage'] = (vc['Count'] / vc['Count'].sum() * 100).round(1)
            fig = px.pie(vc, names=fcat[0], values='Count', hole=0.42,
                         color_discrete_sequence=cs,
                         hover_data={'Percentage': True})
            fig.update_traces(
                textinfo='percent+label',
                textfont_size=12,
                textfont_color=t['text'],
                hovertemplate='<b>%{label}</b><br>Count: %{value:,}<br>Share: %{percent}<extra></extra>'
            )
            fig.update_layout(**base_layout(f'Distribution of {fcat[0]}', t))
            charts.append(make_chart(fig, f'Distribution of {fcat[0]}', '🍩'))
        safe(_pie)

    # ── 3. Time-series line chart (if date column exists) ──
    if date_col and fc:
        def _timeseries():
            tmp = df.copy()
            tmp['_date'] = date_series
            tmp = tmp.dropna(subset=['_date']).sort_values('_date')
            fig = go.Figure()
            for i, col in enumerate(fc[:4]):
                s = clean_numeric(tmp, col)
                if len(s) < 2: continue
                fig.add_trace(go.Scatter(
                    x=tmp['_date'], y=tmp[col],
                    mode='lines+markers', name=col,
                    line=dict(color=cs[i % len(cs)], width=2),
                    marker=dict(size=4),
                    hovertemplate=f'<b>{col}</b><br>Date: %{{x}}<br>Value: %{{y:,.2f}}<extra></extra>'
                ))
            fig.update_layout(**base_layout('Time Series Trend', t))
            charts.append(make_chart(fig, 'Time Series Trend', '📅'))
        safe(_timeseries)
    elif fc:
        # ── 3b. Area trend chart ──
        def _line():
            fig = go.Figure()
            for i, col in enumerate(fc[:4]):
                s = clean_numeric(df, col)
                if len(s) < 2: continue
                fig.add_trace(go.Scatter(
                    y=s.values, mode='lines', name=col,
                    fill='tozeroy' if i == 0 else 'none',
                    line=dict(color=cs[i % len(cs)], width=2.5),
                    fillcolor='rgba(137,180,250,0.06)',
                    hovertemplate=f'<b>{col}</b><br>Index: %{{x}}<br>Value: %{{y:,.2f}}<extra></extra>'
                ))
            fig.update_layout(**base_layout('Trend Analysis', t))
            charts.append(make_chart(fig, 'Trend Analysis', '📈'))
        safe(_line)

    # ── 4. Histogram with stats overlay ──
    if fc:
        def _hist():
            s = clean_numeric(df, fc[0])
            mean_val = s.mean()
            median_val = s.median()
            fig = px.histogram(df, x=fc[0], nbins=30,
                               color_discrete_sequence=[t['accent']],
                               marginal='box',
                               hover_data=df.columns)
            fig.add_vline(x=mean_val, line_dash='dash', line_color=t['accent2'],
                          annotation_text=f'Mean: {mean_val:,.2f}',
                          annotation_font_color=t['accent2'])
            fig.add_vline(x=median_val, line_dash='dot', line_color=t['green'] if 'green' in t else '#a6e3a1',
                          annotation_text=f'Median: {median_val:,.2f}',
                          annotation_font_color=t.get('green','#a6e3a1'),
                          annotation_position='bottom right')
            fig.update_layout(**base_layout(f'Distribution of {fc[0]}', t))
            fig.update_traces(marker_line_width=0)
            charts.append(make_chart(fig, f'Distribution of {fc[0]}', '📉'))
        safe(_hist)

    # ── 5. Scatter with trendline & R² ──
    if len(fc) >= 2:
        def _scatter():
            tmp = df[[fc[0], fc[1]]].dropna()
            if len(tmp) < 5: return
            kw = dict(color=fcat[0], color_discrete_sequence=cs) if fcat else dict(color_discrete_sequence=[t['accent2']])
            try:
                fig = px.scatter(tmp, x=fc[0], y=fc[1], trendline='ols',
                                 trendline_color_override=t['accent2'],
                                 opacity=0.75, **kw)
            except Exception:
                fig = px.scatter(tmp, x=fc[0], y=fc[1],
                                 opacity=0.75, **kw)
            fig.update_traces(marker=dict(size=7, line=dict(width=0)))
            fig.update_layout(**base_layout(f'{fc[0]} vs {fc[1]}', t))
            charts.append(make_chart(fig, f'{fc[0]} vs {fc[1]}', '🔵'))
        safe(_scatter)

    # ── 6. Box plot with outlier detection ──
    if fc and fcat:
        def _box():
            fig = px.box(df, x=fcat[0], y=fc[0],
                         color=fcat[0], color_discrete_sequence=cs,
                         points='outliers', notched=False,
                         hover_data={fc[0]: ':.2f'})
            fig.update_layout(**base_layout(f'{fc[0]} spread by {fcat[0]}', t))
            fig.update_traces(marker_size=4, line_width=1.5)
            charts.append(make_chart(fig, f'{fc[0]} spread by {fcat[0]}', '📦'))
        safe(_box)

    # ── 7. Correlation Heatmap (accurate values) ──
    if len(fc) >= 2:
        def _heatmap():
            corr = df[fc].corr(method='pearson').round(3)
            fig = px.imshow(corr,
                            text_auto='.2f',
                            color_continuous_scale='RdBu_r',
                            zmin=-1, zmax=1, aspect='auto',
                            labels=dict(color='Correlation'))
            fig.update_traces(
                hovertemplate='<b>%{x}</b> vs <b>%{y}</b><br>Correlation: %{z:.3f}<extra></extra>'
            )
            fig.update_layout(**base_layout('Correlation Heatmap', t))
            charts.append(make_chart(fig, 'Correlation Heatmap', '🔥'))
        safe(_heatmap)

    # ── 8. Grouped Bar (2 metrics side by side) ──
    if len(fc) >= 2 and fcat:
        def _gbar():
            g = df.groupby(fcat[0])[fc[:3]].sum().reset_index().head(12)
            fig = px.bar(g, x=fcat[0], y=fc[:3],
                         barmode='group',
                         color_discrete_sequence=cs,
                         text_auto='.2s')
            fig.update_layout(**base_layout(f'Comparison: {" vs ".join(fc[:3])}', t))
            fig.update_traces(marker_line_width=0, textposition='outside')
            charts.append(make_chart(fig, f'Multi-metric Comparison', '📊'))
        safe(_gbar)

    # ── 9. Violin plot (distribution shape) ──
    if fc and fcat and df[fcat[0]].nunique() <= 10:
        def _violin():
            fig = px.violin(df, x=fcat[0], y=fc[0],
                            color=fcat[0], color_discrete_sequence=cs,
                            box=True, points='outliers')
            fig.update_layout(**base_layout(f'{fc[0]} distribution by {fcat[0]}', t))
            charts.append(make_chart(fig, f'{fc[0]} violin', '🎻'))
        safe(_violin)

    # ── 10. Cumulative line chart ──
    if fc:
        def _cumulative():
            s = clean_numeric(df, fc[0]).sort_values().reset_index(drop=True)
            cumsum = s.cumsum()
            fig = go.Figure()
            fig.add_trace(go.Scatter(
                x=list(range(len(s))), y=cumsum,
                mode='lines', name=f'Cumulative {fc[0]}',
                line=dict(color=cs[0], width=2),
                fill='tozeroy', fillcolor=f'rgba(137,180,250,0.07)',
                hovertemplate=f'Index: %{{x}}<br>Cumulative: %{{y:,.2f}}<extra></extra>'
            ))
            fig.update_layout(**base_layout(f'Cumulative {fc[0]}', t))
            charts.append(make_chart(fig, f'Cumulative {fc[0]}', '📐'))
        safe(_cumulative)

    # ── 11. 3D Scatter ──
    if len(fc) >= 3:
        def _3d():
            fig = px.scatter_3d(df.head(300), x=fc[0], y=fc[1], z=fc[2],
                                color=fcat[0] if fcat else None,
                                color_discrete_sequence=cs,
                                opacity=0.75)
            fig.update_layout(
                paper_bgcolor='rgba(0,0,0,0)',
                font=dict(color=t['text']),
                margin=dict(l=0, r=0, t=45, b=0)
            )
            charts.append(make_chart(fig, '3D Analysis', '🌐'))
        safe(_3d)

    # ── 12. Funnel / Waterfall for ranking ──
    if fcat and fc and ('rank' in a or 'top' in a or 'sales' in a or 'revenue' in a):
        def _funnel():
            g = df.groupby(fcat[0])[fc[0]].sum().reset_index()
            g = g.sort_values(fc[0], ascending=False).head(8)
            if len(g) >= 2:
                fig = px.funnel(g, x=fc[0], y=fcat[0],
                                color_discrete_sequence=cs)
                fig.update_layout(**base_layout(f'Funnel: {fc[0]} by {fcat[0]}', t))
                charts.append(make_chart(fig, 'Funnel Ranking', '🔻'))
        safe(_funnel)

    return charts
    # 8. Grouped bar (2 metrics)
    if len(fc)>=2 and fcat:
        def _gbar():
            g=df.groupby(fcat[0])[fc[:2]].sum().reset_index().head(12)
            fig=px.bar(g,x=fcat[0],y=fc[:2],barmode='group',color_discrete_sequence=cs)
            fig.update_layout(**base_layout(f'{fc[0]} vs {fc[1]}',t))
            charts.append(make_chart(fig,f'Comparison: {fc[0]} vs {fc[1]}','📊'))
        safe(_gbar)

    # 9. Violin plot
    if fc and fcat and len(df[fcat[0]].unique())<=12:
        def _violin():
            fig=px.violin(df,x=fcat[0],y=fc[0],color=fcat[0],
                          color_discrete_sequence=cs,box=True)
            fig.update_layout(**base_layout(f'{fc[0]} distribution',t))
            charts.append(make_chart(fig,f'{fc[0]} violin','🎻'))
        safe(_violin)

    # 10. 3D scatter
    if len(fc)>=3:
        def _3d():
            fig=px.scatter_3d(df.head(300),x=fc[0],y=fc[1],z=fc[2],
                              color=fcat[0] if fcat else None,color_discrete_sequence=cs)
            fig.update_layout(paper_bgcolor='rgba(0,0,0,0)',
                              font=dict(color=t['text']),margin=dict(l=0,r=0,t=45,b=0))
            charts.append(make_chart(fig,'3D Analysis','🌐'))
        safe(_3d)

    # 11. Funnel for ranking
    if fcat and fc and ('rank' in a or 'top' in a or 'compar' in a or 'sales' in a):
        def _funnel():
            g=df.groupby(fcat[0])[fc[0]].sum().reset_index().sort_values(fc[0],ascending=False).head(8)
            if len(g)>=2:
                fig=px.funnel(g,x=fc[0],y=fcat[0])
                fig.update_layout(**base_layout('Funnel Ranking',t))
                charts.append(make_chart(fig,'Funnel Ranking','🔻'))
        safe(_funnel)

    # 12. Sunburst hierarchy
    if len(fcat)>=2 and fc:
        def _sun():
            fig=px.sunburst(df.head(400),path=fcat[:2],values=fc[0],
                            color_discrete_sequence=cs)
            fig.update_layout(**base_layout('Hierarchy View',t))
            charts.append(make_chart(fig,'Hierarchy View','☀️'))
        safe(_sun)

    return charts

def build_summary(df):
    num = df.select_dtypes(include='number')
    missing = df.isnull().sum()
    cards = []
    for col in num.columns[:6]:
        s = num[col].dropna()
        if len(s) == 0: continue
        cards.append({
            'label': col,
            'total': f"{s.sum():,.2f}",
            'avg':   f"{s.mean():,.2f}",
            'max':   f"{s.max():,.2f}",
            'min':   f"{s.min():,.2f}",
            'median':f"{s.median():,.2f}",
            'std':   f"{s.std():,.2f}",
            'count': int(s.count()),
        })

    # Detect date column
    date_col, _ = smart_detect_date_col(df)

    # Duplicate count
    duplicate_count = int(df.duplicated().sum())

    return {
        'rows':               int(df.shape[0]),
        'columns':            int(df.shape[1]),
        'column_names':       list(df.columns),
        'numeric_columns':    list(num.columns),
        'categorical_columns':list(df.select_dtypes(include='object').columns),
        'missing_values':     {c: int(v) for c, v in missing.items() if v > 0},
        'total_missing':      int(missing.sum()),
        'duplicate_rows':     duplicate_count,
        'date_column':        date_col,
        'stat_cards':         cards,
        'memory_kb':          round(df.memory_usage(deep=True).sum() / 1024, 1),
    }

def build_prompt(df,filename,analysis_type,questions,specific_cols):
    num=df.select_dtypes(include='number')
    stats=num.describe().round(2).to_string() if not num.empty else "No numeric data"
    sample=df.head(8).to_string()
    qs='\n'.join(f'- {q}' for q in questions) if questions else '- General overview'
    cols=', '.join(specific_cols) if specific_cols else 'All columns'
    return f"""You are a world-class data analyst. Analyse this dataset thoroughly and professionally.

FILE: {filename}
SHAPE: {df.shape[0]:,} rows × {df.shape[1]} columns
COLUMNS: {', '.join(df.columns)}
FOCUS: {cols}
ANALYSIS TYPE: {analysis_type}

USER QUESTIONS:
{qs}

STATISTICS:
{stats}

SAMPLE DATA (first 8 rows):
{sample}

MISSING VALUES: {df.isnull().sum().to_dict()}

Write a professional report with these sections:
## 📋 Executive Summary
## 🔍 Direct Answers to Your Questions
## 📊 Key Findings (use real numbers)
## 💡 Business Insights
## ⚠️ Risks & Anomalies
## ✅ Recommendations (5 specific actions)
## 📈 Data Quality Assessment

Be specific, use actual numbers, use bullet points."""

def fallback_analysis(df,filename):
    num=df.select_dtypes(include='number')
    lines=[f"## 📋 Executive Summary\n**{filename}** — {df.shape[0]:,} rows × {df.shape[1]} columns.\n"]
    lines.append("## 📊 Key Findings\n")
    for col in num.columns[:6]:
        s=num[col].dropna()
        lines.append(f"- **{col}**: Total={s.sum():,.2f} | Avg={s.mean():,.2f} | Min={s.min():,.2f} | Max={s.max():,.2f}")
    for col in df.select_dtypes(include='object').columns[:3]:
        top=df[col].value_counts().head(3)
        lines.append(f"- **{col}** top values: {', '.join(f'{k} ({v})' for k,v in top.items())}")
    missing=df.isnull().sum(); bad=missing[missing>0]
    lines.append("\n## ⚠️ Data Quality\n")
    lines.append("- ✅ No missing values." if bad.empty else
                 '\n'.join(f"- ⚠️ **{c}**: {v} missing ({v/len(df)*100:.1f}%)" for c,v in bad.items()))
    lines.append("\n## ✅ Recommendations\n- Add a Gemini API key in app.py for full AI analysis (it's free!).")
    return '\n'.join(lines)

# ── Auth Routes ───────────────────────────────────────────────────────────────
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/register', methods=['GET','POST'])
def register():
    if current_user.is_authenticated: return redirect(url_for('dashboard'))
    if request.method=='POST':
        username=request.form.get('username','').strip()
        email=request.form.get('email','').strip()
        password=request.form.get('password','')
        confirm=request.form.get('confirm','')
        if not username or not email or not password:
            flash('All fields are required.','error'); return render_template('register.html', users_exist=User.query.count()>0)
        if password != confirm:
            flash('Passwords do not match.','error'); return render_template('register.html', users_exist=User.query.count()>0)
        if len(password)<6:
            flash('Password must be at least 6 characters.','error'); return render_template('register.html', users_exist=User.query.count()>0)
        if User.query.filter_by(email=email).first():
            flash('Email already registered.','error'); return render_template('register.html', users_exist=User.query.count()>0)
        if User.query.filter_by(username=username).first():
            flash('Username already taken.','error'); return render_template('register.html', users_exist=User.query.count()>0)
        role='admin' if User.query.count()==0 else 'user'
        u=User(username=username,email=email,role=role)
        u.set_password(password); db.session.add(u); db.session.commit()
        login_user(u)
        flash(f'Welcome, {username}! Account created.','success')
        return redirect(url_for('dashboard'))
    return render_template('register.html', users_exist=User.query.count()>0)

@app.route('/login', methods=['GET','POST'])
def login():
    if current_user.is_authenticated: return redirect(url_for('dashboard'))
    if request.method=='POST':
        email=request.form.get('email','').strip()
        password=request.form.get('password','')
        u=User.query.filter_by(email=email).first()
        if u and u.check_password(password):
            login_user(u, remember=True)
            flash(f'Welcome back, {u.username}!','success')
            return redirect(url_for('dashboard'))
        flash('Invalid email or password.','error')
    return render_template('login.html')

@app.route('/logout')
@login_required
def logout():
    logout_user(); flash('Logged out successfully.','success')
    return redirect(url_for('index'))

@app.route('/dashboard')
@login_required
def dashboard():
    return render_template('dashboard.html')

@app.route('/analyse')
@login_required
def analyse():
    return render_template('analyse.html')

@app.route('/settings', methods=['GET','POST'])
@login_required
def settings():
    if request.method=='POST':
        action=request.form.get('action')
        if action=='profile':
            un=request.form.get('username','').strip()
            em=request.form.get('email','').strip()
            if un and un != current_user.username:
                if User.query.filter_by(username=un).first():
                    flash('Username taken.','error')
                else:
                    current_user.username=un; db.session.commit(); flash('Username updated.','success')
            if em and em != current_user.email:
                if User.query.filter_by(email=em).first():
                    flash('Email in use.','error')
                else:
                    current_user.email=em; db.session.commit(); flash('Email updated.','success')
        elif action=='password':
            old=request.form.get('old_password','')
            new=request.form.get('new_password','')
            confirm=request.form.get('confirm_password','')
            if not current_user.check_password(old):
                flash('Current password incorrect.','error')
            elif new != confirm:
                flash('New passwords do not match.','error')
            elif len(new)<6:
                flash('Password must be at least 6 chars.','error')
            else:
                current_user.set_password(new); db.session.commit(); flash('Password updated.','success')
        elif action=='theme':
            th=request.form.get('theme','dark')
            if th in THEMES:
                current_user.theme=th; db.session.commit(); flash(f'Theme changed to {th}.','success')
        return redirect(url_for('settings'))
    return render_template('settings.html', themes=THEMES)

@app.route('/admin')
@login_required
def admin():
    if current_user.role != 'admin':
        flash('Access denied.','error'); return redirect(url_for('dashboard'))
    users=User.query.order_by(User.created_at.desc()).all()
    total_analyses=sum(u.analyses_count for u in users)
    return render_template('admin.html', users=users, total_analyses=total_analyses)

@app.route('/admin/delete_user/<int:uid>', methods=['POST'])
@login_required
def delete_user(uid):
    if current_user.role != 'admin':
        return jsonify({'error':'Access denied'}),403
    u=User.query.get_or_404(uid)
    if u.role=='admin':
        flash('Cannot delete admin users.','error')
    else:
        db.session.delete(u); db.session.commit(); flash(f'User {u.username} deleted.','success')
    return redirect(url_for('admin'))

# ── API Routes ────────────────────────────────────────────────────────────────
@app.route('/providers')
def providers():
    return jsonify({p: key_set(p) for p in ['gemini','chatgpt','claude','deepseek']})

@app.route('/upload', methods=['POST'])
@login_required
def upload():
    if 'file' not in request.files: return jsonify({'error':'No file uploaded'}),400
    file=request.files['file']
    if not file.filename or get_ext(file.filename) not in ALLOWED:
        return jsonify({'error':'Unsupported file type'}),400
    try:
        df=read_file(file)
        if df is None or df.empty: return jsonify({'error':'File empty or unreadable'}),400
        for col in df.columns:
            try: df[col]=pd.to_numeric(df[col].astype(str).str.replace(',','').str.strip())
            except: pass
        df.to_json(tmp_path(), orient='records')
        return jsonify({
            'success':True, 'filename':file.filename,
            'filetype':get_ext(file.filename).upper(),
            'rows':int(df.shape[0]), 'columns':int(df.shape[1]),
            'numeric_columns':df.select_dtypes(include='number').columns.tolist(),
            'categorical_columns':df.select_dtypes(include='object').columns.tolist(),
            'all_columns':list(df.columns),
            'preview':df.head(3).fillna('').astype(str).to_dict(orient='records'),
        })
    except Exception as e: return jsonify({'error':str(e)}),500

@app.route('/analyse_data', methods=['POST'])
@login_required
def analyse_data():
    try:
        p=request.json
        analysis_type=p.get('analysis_type','General Overview')
        questions=p.get('questions',[])
        specific_cols=p.get('columns',[])
        filename=p.get('filename','data')
        provider=p.get('provider', get_active_provider() or 'gemini')
        if not os.path.exists(tmp_path()):
            return jsonify({'error':'Session expired. Please re-upload.'}),400
        df=pd.read_json(tmp_path())
        theme=get_theme()
        summary=build_summary(df)
        charts=generate_charts(df, analysis_type, specific_cols or None, theme)
        preview=df.head(15).fillna('').astype(str).to_dict(orient='records')
        if key_set(provider):
            try: ai_text=call_ai(provider, build_prompt(df,filename,analysis_type,questions,specific_cols))
            except Exception as e: ai_text=fallback_analysis(df,filename)+f"\n\n⚠️ AI error: {e}"
        else:
            ai_text=fallback_analysis(df,filename)
        current_user.analyses_count+=1; db.session.commit()
        return jsonify({'success':True,'filename':filename,'analysis_type':analysis_type,
                        'provider':provider,'summary':summary,'charts':charts,
                        'preview':preview,'columns':list(df.columns),'ai_analysis':ai_text})
    except Exception as e: return jsonify({'error':str(e)}),500

@app.route('/chat', methods=['POST'])
@login_required
def chat():
    try:
        data=request.json
        question=data.get('question','')
        context=data.get('context','')
        provider=data.get('provider', get_active_provider() or 'gemini')
        if not key_set(provider):
            names={'gemini':'Gemini','chatgpt':'ChatGPT','claude':'Claude','deepseek':'DeepSeek'}
            return jsonify({'answer':f'⚙️ No API key for **{names.get(provider,provider)}**.\n\nAdd it in app.py → API_KEYS.\n\nFree Gemini key: https://aistudio.google.com'})
        sample=stats=''
        if os.path.exists(tmp_path()):
            df=pd.read_json(tmp_path())
            sample=df.head(10).to_string()
            try: stats=df.select_dtypes(include='number').describe().round(2).to_string()
            except: stats=''
        prompt=f"""You are a professional data analyst for Biz-Analyst.
Dataset: {context}
Stats:\n{stats}
Sample:\n{sample}
Question: {question}
Answer clearly using actual numbers. Be concise but thorough."""
        return jsonify({'answer':call_ai(provider,prompt)})
    except Exception as e: return jsonify({'answer':f'❌ Error: {str(e)}'})

@app.route('/update_theme', methods=['POST'])
@login_required
def update_theme():
    th=request.json.get('theme','dark')
    if th in THEMES:
        current_user.theme=th; db.session.commit()
        return jsonify({'success':True,'theme':THEMES[th]})
    return jsonify({'error':'Invalid theme'}),400

# ── Custom AI Model ───────────────────────────────────────────────────────────

class CustomAI(db.Model):
    id           = db.Column(db.Integer, primary_key=True)
    user_id      = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    name         = db.Column(db.String(80), nullable=False, default='My AI')
    avatar       = db.Column(db.String(10), default='🤖')
    personality  = db.Column(db.Text, default='')
    expertise    = db.Column(db.String(200), default='data analysis')
    knowledge    = db.Column(db.Text, default='')   # uploaded docs text
    provider     = db.Column(db.String(20), default='gemini')
    language     = db.Column(db.String(20), default='English')
    tone         = db.Column(db.String(20), default='professional')
    created_at   = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at   = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

    def build_system_prompt(self):
        tone_desc = {
            'professional': 'formal, professional and precise',
            'friendly':     'warm, friendly and approachable',
            'casual':       'casual, relaxed and conversational',
            'expert':       'highly technical, detailed and expert-level',
            'simple':       'simple, clear and easy to understand for beginners',
        }.get(self.tone, 'professional')

        parts = [
            f"You are {self.name}, a custom AI assistant.",
            f"Your communication style is {tone_desc}.",
            f"Your area of expertise is: {self.expertise}.",
        ]
        if self.personality:
            parts.append(f"Your personality: {self.personality}")
        if self.knowledge:
            parts.append(f"\nYour knowledge base (use this to answer questions):\n{self.knowledge[:3000]}")
        parts.append(f"\nAlways respond in {self.language}.")
        parts.append("If asked about topics outside your expertise, politely redirect to your area of focus.")
        return '\n'.join(parts)


@app.route('/my_ai')
@login_required
def my_ai():
    ai = CustomAI.query.filter_by(user_id=current_user.id).first()
    return render_template('my_ai.html', ai=ai)


@app.route('/my_ai/save', methods=['POST'])
@login_required
def save_ai():
    try:
        data = request.form
        ai = CustomAI.query.filter_by(user_id=current_user.id).first()
        if not ai:
            ai = CustomAI(user_id=current_user.id)
            db.session.add(ai)

        ai.name        = data.get('name', 'My AI').strip() or 'My AI'
        ai.avatar      = data.get('avatar', '🤖').strip() or '🤖'
        ai.personality = data.get('personality', '').strip()
        ai.expertise   = data.get('expertise', 'data analysis').strip()
        ai.provider    = data.get('provider', 'gemini')
        ai.language    = data.get('language', 'English')
        ai.tone        = data.get('tone', 'professional')

        # Handle knowledge file upload
        if 'knowledge_file' in request.files:
            kfile = request.files['knowledge_file']
            if kfile and kfile.filename:
                try:
                    ext = get_ext(kfile.filename)
                    if ext == 'pdf' and PDF_OK:
                        import pdfplumber as _pdf
                        with _pdf.open(io.BytesIO(kfile.read())) as pdf:
                            texts = [p.extract_text() for p in pdf.pages if p.extract_text()]
                        ai.knowledge = '\n'.join(texts)[:8000]
                    elif ext in ('docx','doc') and DOCX_OK:
                        doc = DocxDoc(io.BytesIO(kfile.read()))
                        ai.knowledge = '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])[:8000]
                    elif ext == 'txt':
                        ai.knowledge = kfile.read().decode('utf-8', errors='replace')[:8000]
                    else:
                        ai.knowledge = kfile.read().decode('utf-8', errors='replace')[:8000]
                except Exception as e:
                    flash(f'Could not read knowledge file: {e}', 'error')

        # Manual knowledge text
        manual_knowledge = data.get('knowledge_text', '').strip()
        if manual_knowledge:
            ai.knowledge = (ai.knowledge or '') + '\n' + manual_knowledge

        ai.updated_at = datetime.utcnow()
        db.session.commit()
        flash(f'✅ {ai.name} saved successfully!', 'success')
    except Exception as e:
        flash(f'Error saving AI: {e}', 'error')
    return redirect(url_for('my_ai'))


@app.route('/my_ai/chat', methods=['POST'])
@login_required
def my_ai_chat():
    try:
        data     = request.json
        question = data.get('message', '').strip()
        history  = data.get('history', [])   # [{role, content}]

        if not question:
            return jsonify({'error': 'No message provided'}), 400

        ai = CustomAI.query.filter_by(user_id=current_user.id).first()
        if not ai:
            return jsonify({'reply': "You haven't created your AI yet. Go to **My AI** to set it up first!"})

        provider = ai.provider
        if not key_set(provider):
            return jsonify({'reply': f"⚙️ Your AI uses **{provider}** but no API key is set.\n\nAdd your {provider} key in `app.py` → `API_KEYS`.\n\nFree Gemini key: https://aistudio.google.com"})

        system = ai.build_system_prompt()

        # Build conversation prompt
        history_text = ''
        for h in history[-8:]:   # last 8 exchanges
            role = 'User' if h.get('role') == 'user' else ai.name
            history_text += f"{role}: {h.get('content','')}\n"

        full_prompt = f"""{system}

Previous conversation:
{history_text}
User: {question}
{ai.name}:"""

        reply = call_ai(provider, full_prompt)
        return jsonify({'reply': reply, 'ai_name': ai.name, 'ai_avatar': ai.avatar})

    except Exception as e:
        return jsonify({'reply': f'❌ Error: {str(e)}'})


@app.route('/my_ai/reset_knowledge', methods=['POST'])
@login_required
def reset_ai_knowledge():
    ai = CustomAI.query.filter_by(user_id=current_user.id).first()
    if ai:
        ai.knowledge = ''
        db.session.commit()
        flash('Knowledge base cleared.', 'success')
    return redirect(url_for('my_ai'))

# ── Data Editor Routes ────────────────────────────────────────────────────────

@app.route('/editor')
@login_required
def editor():
    return render_template('editor.html')

@app.route('/editor/load', methods=['POST'])
@login_required
def editor_load():
    """Load data into the editor — from uploaded file or existing session."""
    source = request.json.get('source','session')
    try:
        if source == 'session':
            if not os.path.exists(tmp_path()):
                return jsonify({'error':'No data loaded. Upload a file first on the Analyse page.'}),400
            df = pd.read_json(tmp_path())
        else:
            return jsonify({'error':'Unknown source'}),400

        # Return full table data
        rows = df.fillna('').astype(str).to_dict(orient='records')
        return jsonify({
            'success': True,
            'columns': list(df.columns),
            'rows': rows,
            'shape': [int(df.shape[0]), int(df.shape[1])],
            'dtypes': {c: str(t) for c, t in df.dtypes.items()},
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/editor/upload', methods=['POST'])
@login_required
def editor_upload():
    """Upload a new file directly into the editor."""
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    file = request.files['file']
    if not file.filename or get_ext(file.filename) not in ALLOWED:
        return jsonify({'error': 'Unsupported file type'}), 400
    try:
        df = read_file(file)
        if df is None or df.empty:
            return jsonify({'error': 'File empty or unreadable'}), 400
        # Auto-convert numeric strings
        for col in df.columns:
            try: df[col] = pd.to_numeric(df[col].astype(str).str.replace(',','').str.strip())
            except: pass
        df.to_json(tmp_path(), orient='records')
        rows = df.fillna('').astype(str).to_dict(orient='records')
        return jsonify({
            'success': True,
            'filename': file.filename,
            'columns': list(df.columns),
            'rows': rows,
            'shape': [int(df.shape[0]), int(df.shape[1])],
            'dtypes': {c: str(t) for c, t in df.dtypes.items()},
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/editor/ai_edit', methods=['POST'])
@login_required
def editor_ai_edit():
    """AI edits the data based on plain English instructions."""
    try:
        payload = request.json
        instruction = payload.get('instruction', '').strip()
        provider    = payload.get('provider', get_active_provider() or 'gemini')

        if not instruction:
            return jsonify({'error': 'No instruction provided'}), 400
        if not os.path.exists(tmp_path()):
            return jsonify({'error': 'No data loaded. Upload a file first.'}), 400

        df = pd.read_json(tmp_path())

        # Save backup for undo
        import shutil
        shutil.copy(tmp_path(), tmp_path().replace('.json','_backup.json'))

        if not key_set(provider):
            return jsonify({'error': f'No API key set for {provider}. Add it in app.py → API_KEYS.'}), 400

        # Build a prompt asking AI to return Python pandas code
        cols_info = {c: str(t) for c, t in df.dtypes.items()}
        sample = df.head(5).to_string()

        prompt = f"""You are a Python pandas expert. The user wants to edit a DataFrame.

DataFrame info:
- Shape: {df.shape[0]} rows × {df.shape[1]} columns
- Columns and types: {cols_info}
- Sample (first 5 rows):
{sample}

User instruction: "{instruction}"

Write ONLY a Python code block that modifies the DataFrame called `df`.
Rules:
- The variable is already named `df`
- Do NOT import pandas or read any file
- Do NOT print anything
- Do NOT reassign `df =` to a completely new unrelated object
- Just write the transformation code
- End with `df = df` to confirm the final result
- Wrap code in ```python ... ```

Examples of valid instructions:
- "delete column Age" → df.drop(columns=['Age'], inplace=True)
- "fill missing values with 0" → df.fillna(0, inplace=True)  
- "add column FullName = FirstName + LastName" → df['FullName'] = df['FirstName'] + ' ' + df['LastName']
- "sort by Revenue descending" → df.sort_values('Revenue', ascending=False, inplace=True)
- "rename column Qty to Quantity" → df.rename(columns={{'Qty':'Quantity'}}, inplace=True)
- "filter rows where Sales > 1000" → df = df[df['Sales'] > 1000]
- "convert Price to numeric" → df['Price'] = pd.to_numeric(df['Price'], errors='coerce')
"""

        ai_response = call_ai(provider, prompt)

        # Extract code block from response
        import re
        code_match = re.search(r'```python\s*(.*?)```', ai_response, re.DOTALL)
        if not code_match:
            code_match = re.search(r'```\s*(.*?)```', ai_response, re.DOTALL)

        if not code_match:
            # Try to use the whole response as code if it looks like code
            code = ai_response.strip()
        else:
            code = code_match.group(1).strip()

        # Safety check — block dangerous operations
        banned = ['import os','import sys','open(','exec(','eval(','__import__',
                  'subprocess','shutil','rmdir','remove(','unlink(']
        for b in banned:
            if b in code:
                return jsonify({'error': f'Instruction contains unsafe operation: {b}'}), 400

        # Execute the AI-generated code
        local_vars = {'df': df.copy(), 'pd': pd}
        try:
            exec(code, {'pd': pd, '__builtins__': {}}, local_vars)
            df_new = local_vars.get('df', df)
        except Exception as exec_err:
            return jsonify({
                'error': f'AI generated code that failed to run: {exec_err}',
                'ai_code': code
            }), 400

        if df_new is None or not isinstance(df_new, pd.DataFrame):
            return jsonify({'error': 'AI edit did not return a valid table'}), 400

        # Save the edited data
        df_new.to_json(tmp_path(), orient='records')

        rows = df_new.fillna('').astype(str).to_dict(orient='records')
        return jsonify({
            'success': True,
            'message': f'✅ Done! {instruction}',
            'ai_code': code,
            'columns': list(df_new.columns),
            'rows': rows,
            'shape': [int(df_new.shape[0]), int(df_new.shape[1])],
            'dtypes': {c: str(t) for c, t in df_new.dtypes.items()},
            'changes': f'Was {df.shape[0]}×{df.shape[1]}, now {df_new.shape[0]}×{df_new.shape[1]}',
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/editor/manual_edit', methods=['POST'])
@login_required
def editor_manual_edit():
    """Apply manual cell edits from the table editor."""
    try:
        payload = request.json
        edits   = payload.get('edits', [])   # [{row, col, value}, ...]
        delrows = payload.get('delete_rows', [])
        delcols = payload.get('delete_cols', [])
        addcol  = payload.get('add_column', None)
        newcolname = payload.get('new_column_name', '')
        renaming = payload.get('rename', {})  # {old: new}

        if not os.path.exists(tmp_path()):
            return jsonify({'error': 'No data loaded'}), 400

        df = pd.read_json(tmp_path())

        # Save backup for undo
        import shutil as _shutil
        _shutil.copy(tmp_path(), tmp_path().replace('.json','_backup.json'))

        # Apply cell edits
        for edit in edits:
            r, c, v = edit.get('row'), edit.get('col'), edit.get('value', '')
            if c in df.columns and 0 <= r < len(df):
                try:
                    # Try numeric
                    df.at[r, c] = pd.to_numeric(v)
                except (ValueError, TypeError):
                    df.at[r, c] = v

        # Delete rows
        if delrows:
            df = df.drop(index=delrows, errors='ignore').reset_index(drop=True)

        # Delete columns
        if delcols:
            df = df.drop(columns=[c for c in delcols if c in df.columns], errors='ignore')

        # Add new column
        if addcol and newcolname:
            df[newcolname] = addcol

        # Rename columns
        if renaming:
            df = df.rename(columns=renaming)

        df.to_json(tmp_path(), orient='records')
        rows = df.fillna('').astype(str).to_dict(orient='records')
        return jsonify({
            'success': True,
            'columns': list(df.columns),
            'rows': rows,
            'shape': [int(df.shape[0]), int(df.shape[1])],
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/editor/download', methods=['POST'])
@login_required
def editor_download():
    """Download the edited data as CSV."""
    try:
        from flask import Response
        if not os.path.exists(tmp_path()):
            return jsonify({'error': 'No data loaded'}), 400
        df = pd.read_json(tmp_path())
        csv_data = df.to_csv(index=False)
        return Response(
            csv_data,
            mimetype='text/csv',
            headers={'Content-Disposition': 'attachment;filename=edited_data.csv'}
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/editor/undo', methods=['POST'])
@login_required
def editor_undo():
    """Simple undo — restore from backup."""
    backup = tmp_path().replace('.json', '_backup.json')
    if os.path.exists(backup):
        import shutil
        shutil.copy(backup, tmp_path())
        df = pd.read_json(tmp_path())
        rows = df.fillna('').astype(str).to_dict(orient='records')
        return jsonify({'success': True, 'columns': list(df.columns), 'rows': rows,
                        'shape': [int(df.shape[0]), int(df.shape[1])]})
    return jsonify({'error': 'No backup to undo to'}), 400


if __name__=='__main__':
    with app.app_context():
        db.create_all()
    port = int(os.environ.get('PORT', 8080))
    app.run(debug=False, host='0.0.0.0', port=port)

